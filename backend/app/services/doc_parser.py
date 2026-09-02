"""文档解析服务：将业务文件解析为文本。

PDF 优先使用原生文本层，扫描 PDF 与图片才调用服务端配置的 OCR 适配器。
"""
from __future__ import annotations

import base64
import ipaddress
import json
import socket
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any
from urllib.parse import urlsplit, urlunsplit

import httpcore
import httpx

from ..config import get_settings

_IMAGE_EXTS = {"jpg", "jpeg", "png", "gif", "bmp", "tiff", "tif", "webp"}
_TEXT_EXTS = {"txt", "md", "markdown", "csv", "tsv", "json", "yaml", "yml", "xml", "log"}


@dataclass(frozen=True)
class _OCRTarget:
    scheme: str
    hostname: str
    port: int
    authority: str
    address: str


@dataclass(frozen=True)
class _OCRConfiguration:
    endpoint_url: str
    api_key: str
    engine: str
    language: str
    timeout_seconds: float
    target: _OCRTarget


def _version_pair(value: str) -> tuple[int, int] | None:
    try:
        major, minor, *_rest = str(value).split(".")
        return int(major), int(minor)
    except (TypeError, ValueError):
        return None


def _assert_pinning_runtime_compatibility() -> None:
    if (
        _version_pair(getattr(httpx, "__version__", "")) not in {(0, 27), (0, 28)}
        or _version_pair(getattr(httpcore, "__version__", "")) != (1, 0)
    ):
        raise RuntimeError("OCR 安全固定 IP 所需的 HTTP 运行时版本未经验证")


class _PinnedOCRTransport(httpx.BaseTransport):
    """Connect to one validated IP while retaining the TLS SNI and Host."""

    def __init__(
        self,
        target: _OCRTarget,
        *,
        inner: httpx.BaseTransport | None = None,
    ) -> None:
        _assert_pinning_runtime_compatibility()
        self._target = target
        self._inner = inner or httpx.HTTPTransport(verify=True, trust_env=False)

    def handle_request(self, request: httpx.Request) -> httpx.Response:
        request_host = request.url.raw_host.decode("ascii").casefold()
        request_port = request.url.port or 443
        if (
            request.url.scheme != self._target.scheme
            or request_host != self._target.hostname.casefold()
            or request_port != self._target.port
        ):
            raise httpx.TransportError("OCR transport attempted to leave its pinned origin")

        extensions = dict(request.extensions)
        extensions["sni_hostname"] = self._target.hostname
        headers = request.headers.copy()
        headers["Host"] = self._target.authority
        pinned_request = httpx.Request(
            request.method,
            request.url.copy_with(host=self._target.address),
            headers=headers,
            stream=request.stream,
            extensions=extensions,
        )
        return self._inner.handle_request(pinned_request)

    def close(self) -> None:
        self._inner.close()


def parse_file(path: str | Path, filename: str = "") -> dict[str, Any]:
    """Compatibility wrapper for callers that still own a local path."""
    p = Path(path)
    name = filename or p.name
    try:
        return parse_bytes(p.read_bytes(), name)
    except Exception as exc:  # noqa: BLE001
        return {"status": "error", "text": "", "message": f"解析失败: {exc}"}


def parse_bytes(content: bytes, filename: str) -> dict[str, Any]:
    """Parse one document directly from durable object bytes."""
    if not isinstance(content, bytes):
        return {"status": "error", "text": "", "message": "文件内容必须是字节数据"}
    name = Path(str(filename or "")).name
    ext = Path(name).suffix.lstrip(".").lower()
    try:
        if ext in _IMAGE_EXTS:
            return _parse_image(content, name)
        if ext == "pdf":
            return _parse_pdf(content, name)
        if ext in ("xlsx", "xlsm"):
            return _parse_excel(content)
        if ext == "xls":
            return _parse_xls(content)
        if ext in ("docx",):
            return _parse_docx(content)
        if ext == "pptx":
            return _parse_pptx(content)
        if ext in _TEXT_EXTS:
            return _parse_text(content)
        return {"status": "error", "text": "", "message": f"暂不支持的文件类型: .{ext}"}
    except Exception as exc:  # noqa: BLE001
        return {"status": "error", "text": "", "message": f"解析失败: {exc}"}


# ──────────────────────────────────────────────
# OCR（PDF / 图片）
# ──────────────────────────────────────────────
def _host_allowlist(value: str) -> set[str]:
    return {
        item.strip().rstrip(".").casefold()
        for item in str(value or "").split(",")
        if item.strip()
    }


def _resolve_ocr_target(
    base_url: str,
    *,
    allowed_hosts: set[str],
    private_host_allowlist: set[str],
) -> _OCRTarget | None:
    try:
        normalized = httpx.URL(base_url)
        hostname = normalized.raw_host.decode("ascii").rstrip(".").casefold()
        port = normalized.port or 443
        authority = normalized.netloc.decode("ascii")
    except (UnicodeError, ValueError):
        return None
    if hostname not in allowed_hosts:
        return None

    allow_private = hostname in private_host_allowlist
    try:
        literal = ipaddress.ip_address(hostname)
    except ValueError:
        literal = None
    if literal is not None:
        if not allow_private and not literal.is_global:
            return None
        return _OCRTarget(
            scheme="https",
            hostname=hostname,
            port=port,
            authority=authority,
            address=str(literal),
        )

    try:
        resolved: list[ipaddress.IPv4Address | ipaddress.IPv6Address] = []
        seen: set[str] = set()
        for item in socket.getaddrinfo(hostname, port, type=socket.SOCK_STREAM):
            address = ipaddress.ip_address(str(item[4][0]).split("%", 1)[0])
            if str(address) not in seen:
                seen.add(str(address))
                resolved.append(address)
    except (OSError, ValueError):
        return None
    if not resolved or (
        not allow_private and any(not address.is_global for address in resolved)
    ):
        return None
    return _OCRTarget(
        scheme="https",
        hostname=hostname,
        port=port,
        authority=authority,
        address=str(resolved[0]),
    )


def _ocr_configuration() -> _OCRConfiguration | None:
    settings = get_settings()
    base_url = str(getattr(settings, "ocr_base_url", "") or "").strip()
    endpoint_path = str(getattr(settings, "ocr_endpoint_path", "") or "").strip()
    api_key = str(getattr(settings, "ocr_api_key", "") or "").strip()
    engine = str(getattr(settings, "ocr_engine", "") or "").strip()
    language = str(getattr(settings, "ocr_language", "") or "").strip()
    if not all((base_url, endpoint_path, api_key, engine, language)):
        return None

    try:
        base = urlsplit(base_url)
        endpoint = urlsplit(endpoint_path)
        base_port = base.port
        if (
            base.scheme.lower() != "https"
            or not base.hostname
            or base_port == 0
            or base.username is not None
            or base.password is not None
            or base.query
            or base.fragment
            or endpoint.scheme
            or endpoint.netloc
            or endpoint.query
            or endpoint.fragment
            or not endpoint.path.startswith("/")
            or any(segment in {".", ".."} for segment in endpoint.path.split("/"))
        ):
            return None
        path = f"{base.path.rstrip('/')}{endpoint.path}"
        endpoint_url = urlunsplit((base.scheme, base.netloc, path, "", ""))
        timeout_seconds = float(getattr(settings, "ocr_timeout_seconds", 0.0))
    except (TypeError, ValueError):
        return None

    if not 1.0 <= timeout_seconds <= 600.0:
        return None
    target = _resolve_ocr_target(
        base_url,
        allowed_hosts=_host_allowlist(
            getattr(settings, "ocr_allowed_hosts", "")
        ),
        private_host_allowlist=_host_allowlist(
            getattr(settings, "ocr_private_host_allowlist", "")
        ),
    )
    if target is None:
        return None
    return _OCRConfiguration(
        endpoint_url=endpoint_url,
        api_key=api_key,
        engine=engine,
        language=language,
        timeout_seconds=timeout_seconds,
        target=target,
    )


def _ocr_parse(
    content: bytes,
    filename: str,
    is_image: bool,
    *,
    config: _OCRConfiguration | None = None,
) -> dict[str, Any]:
    config = config or _ocr_configuration()
    if config is None:
        return {"status": "error", "text": "", "message": "OCR 服务未配置"}
    table = True if not is_image else False
    rotate = False if not is_image else True
    try:
        with httpx.Client(
            transport=_PinnedOCRTransport(config.target),
            timeout=config.timeout_seconds,
            verify=True,
            follow_redirects=False,
            trust_env=False,
        ) as client:
            resp = client.post(
                config.endpoint_url,
                files={"file": (filename, content, _mime(filename))},
                data={
                    "backend": config.engine,
                    "lang_list": config.language,
                    "table_enable": str(table).lower(),
                    "auto_rotate": str(rotate).lower(),
                },
                headers={"Authorization": config.api_key},
            )
        if resp.status_code >= 400:
            return {"status": "error", "text": "", "message": "OCR 服务请求失败"}
        try:
            data = resp.json()
        except (TypeError, ValueError):
            return {"status": "error", "text": "", "message": "OCR 服务响应无效"}
        if not isinstance(data, dict) or data.get("status") == "failed":
            return {"status": "error", "text": "", "message": "OCR 服务解析失败"}
        text = data.get("markdown") or data.get("text") or ""
        if not isinstance(text, str):
            return {"status": "error", "text": "", "message": "OCR 服务响应无效"}
        if not text.strip():
            return {"status": "error", "text": "", "message": "OCR 未提取到文本"}
        return {"status": "success", "text": text, "message": "OCR 解析完成"}
    except httpx.TimeoutException:
        return {"status": "error", "text": "", "message": "OCR 服务请求超时"}
    except Exception:  # noqa: BLE001
        return {"status": "error", "text": "", "message": "OCR 服务请求失败"}


def _mime(name: str) -> str:
    import mimetypes

    return mimetypes.guess_type(name)[0] or "application/octet-stream"


def _parse_pdf(content: bytes, name: str) -> dict[str, Any]:
    try:
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(content))
        page_texts: list[str] = []
        for i, page in enumerate(reader.pages):
            page_texts.append(page.extract_text() or "")
        if any(text.strip() for text in page_texts):
            text = "\n\n".join(
                f"--- 第 {i + 1} 页 ---\n{page_text}"
                for i, page_text in enumerate(page_texts)
            )
            return {"status": "success", "text": text, "message": "本地 pypdf 提取"}
    except Exception:  # noqa: BLE001
        pass

    config = _ocr_configuration()
    if config is None:
        return {"status": "error", "text": "", "message": "PDF 无文本层且未配置 OCR 服务"}
    result = _ocr_parse(content, name, is_image=False, config=config)
    if result["status"] == "success":
        return result
    return {
        "status": "error",
        "text": "",
        "message": f"PDF OCR 失败: {result['message']}",
    }


def _parse_image(content: bytes, name: str) -> dict[str, Any]:
    config = _ocr_configuration()
    if config is not None:
        r = _ocr_parse(content, name, is_image=True, config=config)
        if r["status"] == "success":
            return r
        return {"status": "error", "text": "", "message": f"图片 OCR 失败: {r['message']}"}
    return {"status": "error", "text": "", "message": "未配置 OCR 服务，无法解析图片"}


# ──────────────────────────────────────────────
# 表格 / 文档
# ──────────────────────────────────────────────
def _parse_excel(content: bytes) -> dict[str, Any]:
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(content), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        lines = [f"### 工作表: {ws.title}"]
        for row in ws.iter_rows(values_only=True):
            if all(v is None for v in row):
                continue
            lines.append(" | ".join("" if v is None else str(v) for v in row))
        parts.append("\n".join(lines))
    wb.close()
    return {"status": "success", "text": "\n\n".join(parts), "message": "Excel 解析完成"}


def _parse_xls(content: bytes) -> dict[str, Any]:
    try:
        import xlrd  # type: ignore

        book = xlrd.open_workbook(file_contents=content)
        parts = []
        for sh in book.sheets():
            lines = [f"### 工作表: {sh.name}"]
            for r in range(sh.nrows):
                lines.append(" | ".join(str(sh.cell_value(r, c)) for c in range(sh.ncols)))
            parts.append("\n".join(lines))
        return {"status": "success", "text": "\n\n".join(parts), "message": "Excel 解析完成"}
    except ImportError:
        return {"status": "error", "text": "", "message": "缺少 xlrd 依赖，无法解析 .xls（请转换为 .xlsx）"}


def _parse_docx(content: bytes) -> dict[str, Any]:
    from docx import Document

    doc = Document(BytesIO(content))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        parts.append("### 表格")
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return {"status": "success", "text": "\n".join(parts), "message": "Word 解析完成"}


def _parse_pptx(content: bytes) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(BytesIO(content))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"### 幻灯片 {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs)
                    if t.strip():
                        parts.append(t)
    return {"status": "success", "text": "\n".join(parts), "message": "PPT 解析完成"}


def _parse_text(raw: bytes) -> dict[str, Any]:
    for enc in ("utf-8", "gbk", "utf-16"):
        try:
            text = raw.decode(enc)
            break
        except (UnicodeDecodeError, LookupError):
            continue
    else:
        import chardet

        enc = chardet.detect(raw).get("encoding") or "utf-8"
        text = raw.decode(enc, errors="replace")
    return {"status": "success", "text": text, "message": "文本解析完成"}


def parse_base64(b64: str, filename: str) -> dict[str, Any]:
    """解析 Base64 内容（供 skill 调用）。"""
    content = base64.b64decode(b64)
    return parse_bytes(content, filename)
